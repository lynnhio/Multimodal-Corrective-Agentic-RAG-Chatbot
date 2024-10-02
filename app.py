import streamlit as st
from pptx import Presentation
from parsers.pdf_parser import PdfParser

file = st.file_uploader("Upload file", type=["pdf", "pptx"])
if file:
    if(file.name.endswith("pptx")):
        st.header("PPTX")
        prs = Presentation(file)
        # Iterate through the slides
        txt = ""
        for page_number, slide in enumerate(prs.slides):
            # Iterate through the shapes in the slide
            curr_text = ""
            for shape in slide.shapes:
                # Check if the shape has text
                if shape.has_text_frame:
                    # Extract the text
                    st.markdown(shape.text)
                    st.write("")
                    if shape.text not in ["", "\n"]:
                        curr_text += f"\n{shape.text}"

    elif file.name.endswith("pdf"):
        # Example usage:
        pdf_parser = PdfParser(output_folder="extracted_files")
        text_list, documents = pdf_parser.extract_from_pdf(file)
        for x in text_list:
            st.markdown(x)
            st.write("\nNEW PAGE:\n")
                    
            
            