from langchain.schema import Document
from pptx import Presentation
import hashlib
import os

def get_description(image_bytes, context):
    # Placeholder for your actual description logic
    return "dummy description", "yes"

class PPTXParser:
    def __init__(self, output_folder):
        self.output_folder = output_folder
        # Create output folder if it doesn't exist
        if not os.path.exists(self.output_folder):
            os.makedirs(self.output_folder)

        self.id = 0
        self.seen_images = set()
        self.text_list = []
        self.documents = []


    def extract_from_file(self, pptx_file):
        # Open the PPTX file
        prs = Presentation(pptx_file)

        # Extract text from slides
        self._extract_text(prs)
        # Extract images from slides
        self._extract_images(prs)

        return self.text_list, self.documents

    def _extract_text(self, prs):
        # Iterate through the slides to extract text
        for page_number, slide in enumerate(prs.slides):
            curr_text = ""
            for shape in slide.shapes:
                # Check if the shape has text
                if shape.has_text_frame and shape.text not in ["", "\n"]:
                    curr_text += f"\n{shape.text}"
            self.text_list.append(curr_text)

    def _extract_images(self, prs):
        # Iterate through the slides to extract images
        for page_number, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == 13:  # 13 is the type for pictures
                    image = shape.image
                    image_bytes = image.blob
                    image_hash = hashlib.md5(image_bytes).hexdigest()

                    # Check if the image has already been seen
                    if image_hash not in self.seen_images:
                        self.seen_images.add(image_hash)
                        try:
                            description, useful = get_description(image_bytes, self.text_list[page_number])
                            print(f"{useful} at page: {page_number + 1}")
                            if useful == "yes":
                                self.text_list[page_number] += f"\n Image Description:\n{description}"
                                image_filename = os.path.join(self.output_folder, f'page_{page_number + 1}_img_{self.id + 1}.png')
                                doc = Document(
                                    page_content=description,
                                    metadata={
                                        'type': 'image',
                                        'img_path': image_filename
                                    }
                                )
                                self.documents.append(doc)

                                # Save the image
                                with open(image_filename, 'wb') as image_file:
                                    image_file.write(image_bytes)

                                self.id += 1
                        except Exception as e:
                            print(f"Error at page {page_number + 1}: {e}")
                            continue

# Example usage:
# parser = PPTXParser(output_folder="extracted_files")
# text_list, documents = parser.extract_from_pptx(uploaded_file)
